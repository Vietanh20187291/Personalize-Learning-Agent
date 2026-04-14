import os
import re
from groq import Groq
from dotenv import load_dotenv
from langchain_community.document_loaders import PyPDFLoader, Docx2txtLoader, TextLoader
from langchain.schema import Document 
from pptx import Presentation 
from langchain.text_splitter import RecursiveCharacterTextSplitter
from rag.vector_store import get_vector_store

# Tải biến môi trường
load_dotenv()

class CustomPPTXLoader:
    def __init__(self, file_path):
        self.file_path = file_path

    def load(self):
        """Đọc file PPTX và trả về định dạng Document của LangChain"""
        try:
            prs = Presentation(self.file_path)
            full_text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        full_text.append(shape.text)
            
            text_content = "\n".join(full_text)
            return [Document(page_content=text_content, metadata={"source": self.file_path})]
        except Exception as e:
            print(f"❌ Lỗi đọc PPTX: {e}")
            return []

class ContentAgent:
    def __init__(self):
        self.api_key = os.getenv("GROQ_KEY_CONTENT")
        self.client = Groq(api_key=self.api_key) if self.api_key else None
        if not self.api_key:
            print("⚠️ GROQ_KEY_CONTENT chưa được cấu hình, ContentAgent sẽ dùng heuristic fallback.")
        self.model = "llama-3.3-70b-versatile"
        
        self.vector_store = None
        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=1000,
            chunk_overlap=200,
            separators=["\n\n", "\n", " ", ""]
        )
        
        # Danh sách môn học chuẩn
        self.subjects = [
            "Vật lý", "Đại số tuyến tính", "Giải tích", "Tin học đại cương",
            "Chuyên đề giới thiệu ngành CNTT", "Ngôn ngữ lập trình C++",
            "Cấu trúc dữ liệu và giải thuật", "Hệ cơ sở dữ liệu",
            "Kiến trúc máy tính", "Xác suất thống kê", "Toán học tính toán",
            "Mạng máy tính", "PP lập trình hướng đối tượng",
            "Kỹ thuật truyền thông", "Cơ sở hệ điều hành"
        ]

    def _detect_subject(self, text_sample: str, file_name: str):
        """
        Nhận diện môn học thông minh : Kết hợp Tên File + Nội Dung + Mapping từ khóa.
        """
        # Fallback không cần LLM nếu thiếu API key.
        if not self.client:
            combined = f"{file_name} {text_sample}".lower()
            heuristic_map = {
                "Cơ sở Hệ điều hành": ["hệ điều hành", "operating system", "process", "deadlock", "semaphore", "kernel"],
                "Vi xử lý": ["vi xử lý", "processor", "cpu", "assembly", "mạch", "microprocessor"],
                "Cơ sở dữ liệu": ["cơ sở dữ liệu", "database", "sql", "mysql", "postgres", "table", "schema"],
                "Mạng máy tính": ["mạng", "network", "tcp", "udp", "ip", "router", "switch", "osi"],
            }
            for subject_name, keywords in heuristic_map.items():
                if subject_name.lower() in combined or any(keyword in combined for keyword in keywords):
                    return subject_name

            for subject_name in self.subjects:
                if subject_name.lower() in combined:
                    return subject_name
            return "Khác"

        # Tạo danh sách môn học dưới dạng chuỗi để đưa vào prompt
        subjects_str = "\n".join([f"- {s}" for s in self.subjects])

        prompt = f"""
        Bạn là Trợ lý Phân loại Tài liệu Học thuật Đại học.
        Nhiệm vụ: Xác định 01 môn học chính xác nhất cho tài liệu dựa trên tên file và nội dung.

        DANH SÁCH MÔN HỌC CHUẨN (Chỉ được chọn 1 trong số này):
        {subjects_str}

        DỮ LIỆU ĐẦU VÀO:
        - Tên file: "{file_name}"
        - Nội dung trích dẫn: "{text_sample[:3000]}"

        QUY TẮC MAPPING (Suy luận từ khóa):
        1. "C", "C basics", "Lập trình C" -> Chọn "Ngôn ngữ lập trình C++".
        2. "Java", "C#", "OOP", "Hướng đối tượng" -> Chọn "PP lập trình hướng đối tượng".
        3. "SQL", "MySQL", "Database", "CSDL" -> Chọn "Hệ cơ sở dữ liệu".
        4. "IP", "OSI", "TCP/IP", "LAN", "WAN" -> Chọn "Mạng máy tính".
        5. "Đạo hàm", "Tích phân", "Limit" -> Chọn "Giải tích".
        6. "Ma trận", "Định thức", "Matrix" -> Chọn "Đại số tuyến tính".
        7. "Thống kê", "Xác suất", "Biến ngẫu nhiên" -> Chọn "Xác suất thống kê".
        
        YÊU CẦU ĐẦU RA:
        - Chỉ trả về duy nhất tên môn học chính xác giống hệt trong danh sách trên.
        - Nếu tài liệu không thuộc bất kỳ môn nào trong danh sách, trả về "Khác".
        - Không giải thích, không thêm dấu câu.
        """
        
        try:
            chat_completion = self.client.chat.completions.create(
                messages=[{"role": "user", "content": prompt}],
                model=self.model,
                temperature=0.1 # Giữ thấp để AI trả lời chính xác, không sáng tạo lung tung
            )
            detected_name = chat_completion.choices[0].message.content.strip()
            
            # 1. Làm sạch kết quả (xóa dấu chấm, ngoặc kép nếu AI lỡ thêm vào)
            detected_name = detected_name.replace('"', '').replace("'", "").rstrip('.')

            # 2. Kiểm tra khớp chính xác 100%
            for s in self.subjects:
                if s.lower() == detected_name.lower():
                    return s
            
            # 3. Fallback: Kiểm tra khớp một phần (Ví dụ AI trả về "C++" thay vì tên đầy đủ)
            for s in self.subjects:
                if detected_name.lower() in s.lower() or s.lower() in detected_name.lower():
                    return s
                    
            return "Khác"
        except Exception as e:
            print(f"⚠️ Lỗi nhận diện môn học: {e}")
            return "Khác"

    def quick_analyze(self, file_path: str):
        """Phân tích nhanh để gợi ý môn học trên giao diện."""
        try:
            file_name = os.path.basename(file_path)
            loader = self._get_loader(file_path)
            if not loader: return "Khác"

            raw_documents = loader.load()
            if not raw_documents: return "Khác" # Check nếu file rỗng

            text_sample = " ".join([doc.page_content for doc in raw_documents[:5]])
            # Gửi cả tên file và nội dung vào hàm nhận diện thông minh
            return self._detect_subject(text_sample, file_name)
        except Exception as e:
            print(f"⚠️ Lỗi phân tích nhanh: {e}")
            return "Khác"

    def process_file(self, file_path: str, manual_subject: str = None):
        """Quy trình nạp tài liệu vào Vector Store."""
        try:
            file_name = os.path.basename(file_path)
            loader = self._get_loader(file_path)
            if not loader: return False

            raw_documents = loader.load()
            if not raw_documents:
                print("⚠️ File không có nội dung text.")
                return False
            
            # Xác định môn học
            # Ưu tiên môn học được truyền từ lớp học (hỗ trợ cả môn custom ngoài danh sách chuẩn).
            if manual_subject and manual_subject.strip() and manual_subject != "Tự động nhận diện":
                detected_subject = manual_subject.strip()
            else:
                full_text_sample = " ".join([doc.page_content for doc in raw_documents[:10]])
                detected_subject = self._detect_subject(full_text_sample, file_name)

            for doc in raw_documents:
                doc.metadata["subject"] = detected_subject

            chunks = self.text_splitter.split_documents(raw_documents)
            
            if chunks:
                if self.vector_store is None:
                    self.vector_store = get_vector_store()
                self.vector_store.add_documents(chunks)
                print(f"✅ Content Agent: Đã nạp {len(chunks)} đoạn vào môn {detected_subject}")
                return {"success": True, "subject": detected_subject}
            return False

        except Exception as e:
            print(f"❌ Lỗi xử lý file Content Agent: {e}")
            return False

    def _get_loader(self, file_path):
        file_ext = os.path.splitext(file_path)[1].lower()
        if file_ext == ".pdf": return PyPDFLoader(file_path)
        if file_ext == ".docx": return Docx2txtLoader(file_path)
        if file_ext == ".pptx": return CustomPPTXLoader(file_path)
        if file_ext == ".txt": return TextLoader(file_path, encoding="utf-8")
        return None
    
    

# Khởi tạo instance
content_agent = ContentAgent()