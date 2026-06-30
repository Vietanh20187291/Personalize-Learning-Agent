import json
import os
import re
import random
import time
from pathlib import Path
from typing import Dict, List, Optional

from dotenv import load_dotenv
from groq import Groq
from langchain_community.document_loaders import Docx2txtLoader, PyPDFLoader, TextLoader
from pptx import Presentation
from sqlalchemy.orm import Session

from config import settings
from db.models import Document, DocumentPublication, LearnerProfile, LearningRoadmap, StudentDocumentEvaluation, Subject, User
from llm_trace import log_llm_error, log_llm_request, log_llm_response
from memory.conversation_memory import get_conversation_memory
from rag.vector_store import get_vector_store

# Tải biến môi trường
load_dotenv()

MATERIAL_BRIEF_CACHE: Dict[str, dict] = {}

class AdaptiveAgent:
    def __init__(self, db: Session):
        self.db = db
        self.api_key = self._resolve_groq_api_key()
        self.request_timeout_seconds = float(getattr(settings, "ADAPTIVE_AGENT_TIMEOUT_SECONDS", 18) or 18)
        api_key_low = self.api_key.lower()
        if self.api_key and "dummy" not in api_key_low and "testing" not in api_key_low and "placeholder" not in api_key_low:
            try:
                self.client = Groq(api_key=self.api_key, timeout=self.request_timeout_seconds)
            except TypeError:
                self.client = Groq(api_key=self.api_key)
        else:
            self.client = None
        self.model = "llama-3.3-70b-versatile"
        
        # Kết nối tới Vector Database (ChromaDB). Nếu embedding lỗi tạm thời,
        # agent vẫn chạy ở chế độ fallback để không làm gián đoạn trải nghiệm học.
        try:
            self.vector_store = get_vector_store()
        except Exception as exc:
            print(f"⚠️ AdaptiveAgent fallback mode (vector store unavailable): {exc}")
            self.vector_store = None

    def _resolve_groq_api_key(self) -> str:
        candidate_names = [
            "GROQ_KEY_ADAPTIVE",
            "GROQ_API_KEY",
            "GROQ_KEY_DEBUG",
            "GROQ_KEY_ASSESSMENT",
        ]
        blocked_tokens = ("dummy", "testing", "placeholder")

        for env_name in candidate_names:
            value = (os.getenv(env_name) or "").strip()
            if not value:
                continue
            lowered = value.lower()
            if any(token in lowered for token in blocked_tokens):
                continue
            return value

        return (os.getenv("GROQ_KEY_ADAPTIVE") or "").strip()

    def _resolve_subject(self, subject: str) -> Subject:
        subject_name = (subject or "").strip()
        if not subject_name:
            raise ValueError("Subject không được rỗng")

        subject_obj = self.db.query(Subject).filter(Subject.name.ilike(subject_name)).first()
        if subject_obj:
            return subject_obj

        subject_obj = Subject(name=subject_name, description=f"Môn {subject_name}")
        self.db.add(subject_obj)
        self.db.flush()
        return subject_obj

    def _build_fallback_roadmap(self, subject_name: str, level: str):
        level_text = (level or "Beginner").upper()
        sessions = [
            {"session": 1, "topic": f"Nền tảng cốt lõi của {subject_name}", "description": "Làm quen mục tiêu học tập và khái niệm nền tảng quan trọng.", "focus_level": level_text},
            {"session": 2, "topic": "Thuật ngữ và mô hình cơ bản", "description": "Hiểu đúng thuật ngữ, mô hình tư duy và cách áp dụng trong bài tập.", "focus_level": level_text},
            {"session": 3, "topic": "Quy trình thực hành chuẩn", "description": "Rèn quy trình giải quyết bài toán theo từng bước rõ ràng.", "focus_level": level_text},
            {"session": 4, "topic": "Phân tích tình huống điển hình", "description": "Áp dụng kiến thức vào tình huống thực tế ở mức độ vừa phải.", "focus_level": level_text},
            {"session": 5, "topic": "Lỗi thường gặp và cách tránh", "description": "Nhận diện sai sót phổ biến và xây dựng chiến lược khắc phục.", "focus_level": level_text},
            {"session": 6, "topic": "Luyện tập theo nhóm kỹ năng", "description": "Củng cố năng lực qua bài tập phân nhóm từ dễ đến khó.", "focus_level": level_text},
            {"session": 7, "topic": "Vận dụng liên kết kiến thức", "description": "Kết nối các mảng kiến thức để giải bài toán tổng hợp.", "focus_level": level_text},
            {"session": 8, "topic": "Kịch bản ứng dụng thực tế", "description": "Mô phỏng bối cảnh thực tiễn và đưa ra phương án xử lý phù hợp.", "focus_level": level_text},
            {"session": 9, "topic": "Tối ưu cách giải", "description": "So sánh nhiều hướng tiếp cận và chọn cách làm hiệu quả hơn.", "focus_level": level_text},
            {"session": 10, "topic": "Ôn tập chiến lược", "description": "Hệ thống hóa toàn bộ nội dung trước khi kiểm tra cuối khóa.", "focus_level": level_text},
            {"session": 11, "topic": "KIỂM TRA TỔNG HỢP CUỐI KHÓA", "description": "Bài thi đánh giá toàn diện mức độ nắm vững kiến thức.", "focus_level": level_text},
        ]
        return sessions

    def _build_document_driven_roadmap(
        self,
        user_id: int,
        subject_id: int,
        subject_name: str,
        current_level: str,
        allowed_filenames: Optional[list] = None,
    ) -> List[dict]:
        user = self.db.query(User).filter(User.id == user_id).first()
        enrolled_class_ids = [c.id for c in (getattr(user, "enrolled_classes", []) or [])]

        query = self.db.query(Document).join(
            DocumentPublication,
            DocumentPublication.doc_id == Document.id,
        ).filter(
            Document.subject_id == subject_id,
            DocumentPublication.is_visible_to_students == True,
        )

        if enrolled_class_ids:
            query = query.filter(Document.class_id.in_(enrolled_class_ids))
        if allowed_filenames:
            query = query.filter(Document.filename.in_(allowed_filenames))

        docs = query.order_by(Document.upload_time.asc(), Document.id.asc()).all()
        if not docs:
            docs = self.db.query(Document).filter(
                Document.subject_id == subject_id,
            ).order_by(Document.upload_time.asc(), Document.id.asc()).all()
            if allowed_filenames:
                docs = [doc for doc in docs if (doc.filename or "") in set(allowed_filenames)]

        eval_rows = self.db.query(StudentDocumentEvaluation).filter(
            StudentDocumentEvaluation.user_id == user_id,
            StudentDocumentEvaluation.subject_id == subject_id,
        ).all()
        eval_map = {int(row.document_id): row for row in eval_rows}

        prioritized = []
        completed = []
        for doc in docs:
            evaluation = eval_map.get(int(doc.id))
            attempts = int(evaluation.attempts or 0) if evaluation else 0
            latest_score = float(evaluation.latest_score) if evaluation and evaluation.latest_score is not None else None
            title = self._clean_text(doc.title or doc.filename or f"Tài liệu {doc.id}")

            if attempts <= 0:
                prioritized.append({
                    "doc": doc,
                    "title": title,
                    "score": -1.0,
                    "priority": 0,
                    "reason": "Chưa học tài liệu này, nên ưu tiên trước.",
                })
                continue

            if latest_score is not None and latest_score < 40.0:
                prioritized.append({
                    "doc": doc,
                    "title": title,
                    "score": latest_score,
                    "priority": 1,
                    "reason": f"Điểm gần nhất {latest_score:.1f} dưới 40, cần học lại kỹ.",
                })
                continue

            completed.append({
                "doc": doc,
                "title": title,
                "score": latest_score if latest_score is not None else 100.0,
                "priority": 2,
                "reason": "Đã học ổn, giữ làm tài liệu ôn tập sau cùng.",
            })

        if not prioritized:
            prioritized = completed[:]

        prioritized.sort(
            key=lambda item: (
                int(item["priority"]),
                float(item["score"]),
                str(item["title"]).lower(),
                str(item["doc"].filename or "").lower(),
            )
        )

        day_pattern = (3, 4)
        roadmap: List[dict] = []
        for idx, item in enumerate(prioritized, start=1):
            block_days = day_pattern[(idx - 1) % len(day_pattern)]
            title = str(item["title"])
            description = (
                f"Tập trung {block_days} ngày cho tài liệu '{title}'. "
                f"{item['reason']}"
            )
            roadmap.append(
                {
                    "session": idx,
                    "topic": title,
                    "description": description[:220],
                    "focus_level": (current_level or "Beginner").upper(),
                    "document_id": int(item["doc"].id),
                    "source_file": str(item["doc"].filename or ""),
                }
            )

        return roadmap or self._build_fallback_roadmap(subject_name, current_level)

    def _build_fallback_tutor_reply(self, subject: str, user_message: str, roadmap_context: str, context_docs: str) -> str:
        short_context = ""
        if context_docs:
            compact = " ".join(context_docs.split())
            short_context = compact[:260]

        return (
            f"Mình đang hỗ trợ bạn ở chế độ dự phòng cho môn {subject}. "
            f"Theo buổi học hiện tại: {roadmap_context[:140]}.\n\n"
            f"Gợi ý nhanh cho câu hỏi của bạn: \"{user_message[:180]}\"\n"
            f"1. Xác định đúng khái niệm/chủ điểm cần trả lời.\n"
            f"2. Liên hệ với ví dụ gần nhất trong bài học rồi so sánh điểm giống và khác.\n"
            f"3. Tự kiểm tra lại bằng 1 câu ngắn: vì sao cách hiểu này đúng trong ngữ cảnh hiện tại?\n\n"
            f"{('Trích từ học liệu: ' + short_context) if short_context else 'Học liệu đang được cập nhật, bạn có thể gửi câu hỏi cụ thể hơn để mình hướng dẫn từng bước.'}"
        )

    def _clean_text(self, text: str):
        return re.sub(r"\s+", " ", text or "").strip()

    def _looks_like_boilerplate(self, line: str):
        low = self._clean_text(line).lower()
        if not low or len(low) < 18:
            return True

        admin_patterns = [
            r"\bgiảng viên\b",
            r"\bgiáo viên\b",
            r"\bemail\b",
            r"\be-mail\b",
            r"\bsđt\b",
            r"\bsdt\b",
            r"\bđiện thoại\b",
            r"\bhotline\b",
            r"\bliên hệ\b",
            r"\bquy định\b",
            r"\bquy chế\b",
            r"\bđiểm danh\b",
            r"\bnộp bài\b",
            r"\bhọc phí\b",
            r"\blịch học\b",
            r"\bhọc phần\b",
        ]
        if any(re.search(pattern, low) for pattern in admin_patterns):
            return True

        if re.search(r"(?:email|e-mail|sđt|sdt|điện thoại)\s*[:：]", low):
            return True

        return False

    def _load_document_text(self, file_path: str):
        if not file_path or not os.path.exists(file_path):
            return ""

        ext = Path(file_path).suffix.lower()
        try:
            if ext == ".pdf":
                docs = PyPDFLoader(file_path).load()
                return "\n\n".join([doc.page_content for doc in docs if (doc.page_content or "").strip()])

            if ext == ".docx":
                docs = Docx2txtLoader(file_path).load()
                return "\n\n".join([doc.page_content for doc in docs if (doc.page_content or "").strip()])

            if ext == ".txt":
                docs = TextLoader(file_path, encoding="utf-8").load()
                return "\n\n".join([doc.page_content for doc in docs if (doc.page_content or "").strip()])

            if ext == ".pptx":
                prs = Presentation(file_path)
                full_text = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        try:
                            if hasattr(shape, "text") and isinstance(shape.text, str) and shape.text.strip():
                                full_text.append(shape.text)
                        except Exception:
                            continue
                return "\n\n".join(full_text)
        except Exception as exc:
            print(f"⚠️ Không đọc được file nguồn {file_path}: {exc}")

        return ""

    def _fetch_document(self, subject: str, source_file: str = "", document_id: Optional[int] = None):
        query = self.db.query(Document)
        subject_name = self._clean_text(subject)

        if document_id:
            doc = query.filter(Document.id == document_id).first()
            if doc:
                return doc

        if source_file:
            doc = query.filter(Document.filename == source_file).first()
            if doc:
                return doc

        if subject_name:
            doc = query.filter(Document.subject.ilike(subject_name)).order_by(Document.upload_time.desc()).first()
            if doc:
                return doc

        return None

    def _resolve_document_file_path(self, doc: Optional[Document]) -> str:
        if not doc:
            return ""

        project_root = Path(__file__).resolve().parents[2]
        backend_root = Path(__file__).resolve().parents[1]
        candidates: List[Path] = []

        raw_path = self._clean_text(getattr(doc, "file_path", ""))
        if raw_path:
            path_obj = Path(raw_path)
            candidates.append(path_obj)
            if not path_obj.is_absolute():
                candidates.append((project_root / raw_path).resolve())
                candidates.append((backend_root / raw_path).resolve())

        filename = self._clean_text(getattr(doc, "filename", ""))
        if filename:
            candidates.append((project_root / "temp_uploads" / filename).resolve())
            candidates.append((backend_root / "temp_uploads" / filename).resolve())

        seen = set()
        for candidate in candidates:
            candidate_str = str(candidate)
            if candidate_str in seen:
                continue
            seen.add(candidate_str)
            try:
                if candidate.exists() and candidate.is_file():
                    return candidate_str
            except Exception:
                continue

        return ""

    def _get_document_chat_context(
        self,
        subject: str,
        source_file: str = "",
        allowed_filenames: list = None,
        document_id: Optional[int] = None,
    ) -> Dict[str, str]:
        doc = self._fetch_document(subject, source_file=source_file, document_id=document_id)
        effective_source = self._clean_text(getattr(doc, "filename", "")) or self._clean_text(source_file)

        raw_text = ""
        resolved_path = self._resolve_document_file_path(doc)
        if resolved_path:
            raw_text = self._load_document_text(resolved_path)

        if not raw_text.strip():
            raw_text = self._collect_rag_text(
                subject,
                source_file=effective_source,
                allowed_filenames=allowed_filenames,
                document_id=document_id,
            )

        normalized_text = self._normalize_material_text(raw_text)
        final_text = normalized_text or self._clean_text(raw_text)

        return {
            "source_file": effective_source or self._clean_text(source_file),
            "content": final_text[:22000],
        }

    def _collect_rag_text(self, subject: str, source_file: str = "", allowed_filenames: list = None, document_id: Optional[int] = None):
        if self.vector_store is None:
            return ""

        search_filter = {"subject": {"$eq": subject}}
        if allowed_filenames:
            search_filter = {"$and": [{"subject": {"$eq": subject}}, {"source": {"$in": allowed_filenames}}]}
        if source_file:
            search_filter = {"$and": [{"subject": {"$eq": subject}}, {"source": {"$eq": source_file}}]}

        if document_id:
            doc = self._fetch_document(subject, source_file=source_file, document_id=document_id)
            if doc and doc.filename:
                search_filter = {"$and": [{"subject": {"$eq": subject}}, {"source": {"$eq": doc.filename}}]}

        query = f"Toàn bộ nội dung tài liệu, bài học và ví dụ của môn {subject}. {source_file}".strip()
        try:
            docs = self.vector_store.similarity_search(query, k=60, filter=search_filter)
            if not docs and source_file:
                docs_all = self.vector_store.similarity_search(query, k=90, filter={"subject": {"$eq": subject}})
                docs = [d for d in docs_all if os.path.basename(str((d.metadata or {}).get("source", ""))) == source_file]
        except Exception:
            docs = []

        chunks = []
        seen = set()
        for doc in docs:
            text = self._clean_text(getattr(doc, "page_content", ""))
            if len(text) < 20:
                continue
            key = text.lower()
            if key in seen:
                continue
            seen.add(key)
            chunks.append(text)
        return "\n".join(chunks)

    def _normalize_material_text(self, raw_text: str):
        lines = []
        seen = set()
        for raw_line in (raw_text or "").splitlines():
            line = self._clean_text(raw_line)
            if not line or self._looks_like_boilerplate(line):
                continue
            if len(line) < 18:
                continue
            key = line.lower()
            if key in seen:
                continue
            seen.add(key)
            lines.append(line)

        if not lines:
            return ""

        return "\n".join(lines)[:28000]

    def _extract_keywords(self, text: str, limit: int = 8):
        tokens = re.findall(r"[A-Za-zÀ-ỹ0-9]{4,}", (text or "").lower())
        stop = {
            "và", "của", "theo", "cho", "trong", "được", "này", "khi", "một", "các", "những",
            "that", "this", "with", "from", "have", "will", "about", "trên", "dưới",
        }
        freq: Dict[str, int] = {}
        for token in tokens:
            if token in stop or token.isdigit():
                continue
            freq[token] = freq.get(token, 0) + 1
        ranked = sorted(freq.items(), key=lambda item: item[1], reverse=True)
        return [item[0] for item in ranked[:limit]]

    def _build_rule_based_material_brief(self, subject: str, source_file: str, cleaned_text: str):
        sentences = [s.strip() for s in re.split(r"(?<=[\.\!\?;:])\s+", cleaned_text or "") if s.strip()]
        key_points = []
        for sentence in sentences:
            if len(sentence) < 24:
                continue
            if sentence not in key_points:
                key_points.append(sentence[:220])
            if len(key_points) == 6:
                break

        if not key_points:
            key_points = self._extract_keywords(cleaned_text, limit=6)

        important_notes = []
        note_patterns = [r"lưu ý", r"chú ý", r"quan trọng", r"cần nhớ", r"đặc biệt", r"không được", r"phải"]
        for sentence in sentences:
            low = sentence.lower()
            if any(pattern in low for pattern in note_patterns):
                candidate = sentence[:220]
                if candidate not in important_notes:
                    important_notes.append(candidate)
            if len(important_notes) == 5:
                break

        if not important_notes:
            important_notes = [
                "Bám vào khái niệm cốt lõi thay vì chỉ học thuộc từng dòng.",
                "Ưu tiên hiểu các ví dụ minh hoạ rồi mới chuyển sang bài tập.",
                "Ghi nhớ các công thức, điều kiện áp dụng và ngoại lệ nếu có.",
            ]

        keywords = self._extract_keywords(cleaned_text, limit=5)
        topic_label = ", ".join(keywords[:3]) if keywords else self._clean_text(source_file) or subject
        summary_seed = " ".join(sentences[:6])[:900]
        if not summary_seed:
            summary_seed = f"Tài liệu {source_file or subject} tập trung vào {topic_label}."

        key_points_text = "\n- ".join([item[:220] for item in key_points[:5]]) if key_points else "chưa trích được ý chính rõ ràng từ tài liệu"
        notes_text = "\n- ".join([item[:220] for item in important_notes[:5]])

        rewritten_material = "\n\n".join([
            f"Tài liệu đã được biên tập lại cho môn {subject}.",
            f"Trọng tâm chính: {topic_label}.",
            summary_seed,
            f"Các ý cần nhớ:\n- {key_points_text}",
            f"Lưu ý quan trọng:\n- {notes_text}",
        ])

        suggested_prompts = [
            f"Hãy giải thích lại phần cốt lõi của {subject} theo ngôn ngữ dễ hiểu.",
            f"Từ tài liệu này, đâu là 3 ý quan trọng nhất cần nhớ?",
            f"Đặt cho mình 3 câu hỏi kiểm tra nhanh từ tài liệu {source_file or subject}.",
            f"Phần nào trong tài liệu này dễ nhầm nhất và vì sao?",
        ]

        return {
            "subject": subject,
            "source_file": source_file,
            "rewritten_title": self._clean_text(source_file) or f"{subject} - Học liệu đã biên tập",
            "rewritten_material": rewritten_material[:5000],
            "summary": summary_seed[:1200],
            "key_points": key_points[:6],
            "important_notes": important_notes[:5],
            "suggested_prompts": suggested_prompts,
        }

    def _rewrite_material_brief(self, subject: str, source_file: str, cleaned_text: str):
        if not self.client:
            return self._build_rule_based_material_brief(subject, source_file, cleaned_text)

        prompt = f"""
Bạn là biên tập viên học liệu cho sinh viên.

NHIỆM VỤ:
- Viết lại tài liệu theo ngôn ngữ dễ hiểu, mạch lạc, ưu tiên góc nhìn học tập.
- Loại bỏ toàn bộ thông tin thừa: tên giảng viên, giới thiệu giảng viên, sđt, email, quy định môn học, số buổi, % điểm danh, nội quy hành chính, lịch học.
- Không bịa thêm kiến thức ngoài nội dung được cung cấp.
- Không nhắc rằng đây là bản tóm tắt hay bản trích xuất.

ĐẦU RA BẮT BUỘC JSON HỢP LỆ:
{{
  "rewritten_title": "...",
  "rewritten_material": "...",
  "summary": "...",
  "key_points": ["...", "...", "..."],
  "important_notes": ["...", "...", "..."],
  "suggested_prompts": ["...", "...", "...", "..."]
}}

QUY TẮC VIẾT:
- rewritten_material phải có cấu trúc ngắn gọn, rõ nghĩa, chia đoạn hoặc gạch đầu dòng.
- summary phải súc tích và bao quát nội dung chính.
- key_points và important_notes phải bám sát tài liệu, ưu tiên phần cần nhớ khi học.
- suggested_prompts là câu hỏi học tập để sinh viên hỏi tiếp gia sư AI.

MÔN: {subject}
TÀI LIỆU: {source_file}

NỘI DUNG ĐÃ LÀM SẠCH:
{cleaned_text[:22000]}
""".strip()

        try:
            started = time.perf_counter()
            log_llm_request("groq", self.model, prompt=prompt, system_prompt="Output valid JSON only.")
            completion = self.client.chat.completions.create(
                messages=[
                    {"role": "system", "content": "Output valid JSON only."},
                    {"role": "user", "content": prompt},
                ],
                model=self.model,
                temperature=0.2,
                max_tokens=2200,
                response_format={"type": "json_object"},
            )
            raw_content = completion.choices[0].message.content
            duration_ms = (time.perf_counter() - started) * 1000.0
            log_llm_response("groq", self.model, response=str(raw_content), duration_ms=duration_ms)
            data = json.loads(raw_content)
            summary = self._clean_text(str(data.get("summary", "")))
            rewritten_material = self._clean_text(str(data.get("rewritten_material", "")))
            rewritten_title = self._clean_text(str(data.get("rewritten_title", ""))) or self._clean_text(source_file) or f"{subject} - Học liệu đã biên tập"

            key_points = data.get("key_points") or []
            if not isinstance(key_points, list):
                key_points = []
            key_points = [self._clean_text(str(item)) for item in key_points if self._clean_text(str(item))][:6]

            important_notes = data.get("important_notes") or []
            if not isinstance(important_notes, list):
                important_notes = []
            important_notes = [self._clean_text(str(item)) for item in important_notes if self._clean_text(str(item))][:5]

            suggested_prompts = data.get("suggested_prompts") or []
            if not isinstance(suggested_prompts, list):
                suggested_prompts = []
            suggested_prompts = [self._clean_text(str(item)) for item in suggested_prompts if self._clean_text(str(item))][:4]

            if not summary or not rewritten_material:
                return self._build_rule_based_material_brief(subject, source_file, cleaned_text)

            return {
                "subject": subject,
                "source_file": source_file,
                "rewritten_title": rewritten_title,
                "rewritten_material": rewritten_material[:5000],
                "summary": summary[:1200],
                "key_points": key_points,
                "important_notes": important_notes,
                "suggested_prompts": suggested_prompts or self._build_rule_based_material_brief(subject, source_file, cleaned_text)["suggested_prompts"],
            }
        except Exception as exc:
            log_llm_error("groq", self.model, error_message=str(exc), duration_ms=0.0)
            print(f"⚠️ Lỗi biên tập học liệu: {exc}")
            return self._build_rule_based_material_brief(subject, source_file, cleaned_text)

    def _get_material_brief(self, subject: str, source_file: str = "", allowed_filenames: list = None, session_topic: str = "", document_id: Optional[int] = None):
        cache_key = f"{(subject or '').strip().lower()}|{document_id or ''}|{(source_file or '').strip().lower()}|{','.join(sorted(allowed_filenames or []))}"
        cached = MATERIAL_BRIEF_CACHE.get(cache_key)
        if cached:
            return cached

        doc = self._fetch_document(subject, source_file=source_file, document_id=document_id)
        file_text = ""
        if doc and doc.file_path:
            file_text = self._load_document_text(doc.file_path)

        effective_source = doc.filename if doc and doc.filename else source_file
        rag_text = ""
        # Avoid Chroma similarity_search here when the opened document file is
        # already available. This keeps document-grounded tutoring stable on
        # Windows where the Rust Chroma client can crash the process.
        if not file_text.strip():
            rag_text = self._collect_rag_text(
                subject,
                source_file=effective_source,
                allowed_filenames=allowed_filenames,
                document_id=document_id,
            )

        combined_parts = [part for part in [file_text, rag_text] if part.strip()]
        combined_text = "\n\n".join(combined_parts)
        cleaned_text = self._normalize_material_text(combined_text)
        if not cleaned_text:
            fallback = {
                "subject": subject,
                "source_file": source_file,
                "rewritten_title": self._clean_text(source_file) or f"{subject} - Học liệu đã biên tập",
                "rewritten_material": f"Mình chưa trích được đủ nội dung từ tài liệu {source_file or subject}. Hãy mở lại file hoặc chọn đúng tài liệu để mình tóm tắt tiếp.",
                "summary": f"Chưa đủ dữ liệu để biên tập tài liệu {source_file or subject}.",
                "key_points": [],
                "important_notes": ["Hãy thử tải lại tài liệu hoặc chọn đúng file đang học."],
                "suggested_prompts": [
                    "Tóm tắt ngắn tài liệu này",
                    "Nêu các ý cần nhớ trong tài liệu",
                    "Giải thích từng phần theo ngôn ngữ dễ hiểu",
                    "Đặt câu hỏi kiểm tra nhanh từ tài liệu",
                ],
            }
            MATERIAL_BRIEF_CACHE[cache_key] = fallback
            return fallback

        brief = self._rewrite_material_brief(subject, effective_source or source_file, cleaned_text)
        brief["session_topic"] = session_topic
        brief["document_id"] = document_id
        MATERIAL_BRIEF_CACHE[cache_key] = brief
        return brief

    def summarize_material(self, subject: str, source_file: str, session_topic: str = "", allowed_filenames: list = None, document_id: Optional[int] = None):
        brief = self._get_material_brief(
            subject,
            source_file=source_file,
            allowed_filenames=allowed_filenames,
            session_topic=session_topic,
            document_id=document_id,
        )
        result = dict(brief)
        result["suggested_prompts"] = brief.get("suggested_prompts") or [
            "Từ tài liệu này, phần nào dễ nhầm nhất?",
            "Hãy tóm tắt lại theo 5 ý chính ngắn gọn",
            "Tạo 5 câu trắc nghiệm bám sát nội dung tài liệu",
        ]
        return result

    # ==========================================
    # 1. HÀM SINH LỘ TRÌNH HỌC (ROADMAP)
    # ==========================================
    def generate_overall_roadmap(self, user_id: int, subject: str, allowed_filenames: list = None, force_level: str = None):
        subject_obj = self._resolve_subject(subject)
        subject_id = subject_obj.id
        subject_name = subject_obj.name

        current_level = force_level
        if not current_level:
            profile = self.db.query(LearnerProfile).filter_by(user_id=user_id, subject_id=subject_id).first()
            if not profile:
                profile = self.db.query(LearnerProfile).filter_by(user_id=user_id, subject=subject_name).first()
            current_level = profile.current_level if profile else "Beginner"

        try:
            final_roadmap = self._build_document_driven_roadmap(
                user_id=user_id,
                subject_id=subject_id,
                subject_name=subject_name,
                current_level=current_level,
                allowed_filenames=allowed_filenames,
            )

            self.db.query(LearningRoadmap).filter_by(user_id=user_id, subject_id=subject_id).delete()

            new_roadmap = LearningRoadmap(
                user_id=user_id,
                subject_id=subject_id,
                subject=subject_name,
                level_assigned=current_level,
                roadmap_data=final_roadmap,
                current_session=1
            )
            self.db.add(new_roadmap)
            self.db.commit()

            return final_roadmap
        except Exception as e:
            print(f"❌ Lỗi sinh lộ trình: {e}")
            self.db.rollback()
            try:
                fallback = self._build_fallback_roadmap(subject_name, current_level)
                self.db.query(LearningRoadmap).filter_by(user_id=user_id, subject_id=subject_id).delete()
                self.db.add(
                    LearningRoadmap(
                        user_id=user_id,
                        subject_id=subject_id,
                        subject=subject_name,
                        level_assigned=current_level,
                        roadmap_data=fallback,
                        current_session=1,
                    )
                )
                self.db.commit()
                return fallback
            except Exception as db_err:
                print(f"❌ Fallback roadmap cũng lỗi: {db_err}")
                self.db.rollback()
                return []

    # ==========================================
    # BUILD PERSONALIZED CONTEXT (Hồ sơ cá nhân hóa)
    # ==========================================
    def _build_student_context(self, user_id: int, subject: str) -> str:
        """Lấy hồ sơ học tập từ DB để inject vào system prompt Tutor Agent."""
        parts: List[str] = []

        # 1. Learner Profile — mức năng lực
        subject_obj = self.db.query(Subject).filter(Subject.name.ilike(subject)).first()
        subject_id = subject_obj.id if subject_obj else None

        level = "Beginner"
        if subject_id:
            profile = self.db.query(LearnerProfile).filter_by(user_id=user_id, subject_id=subject_id).first()
            if not profile:
                profile = self.db.query(LearnerProfile).filter_by(user_id=user_id, subject=subject).first()
            if profile and profile.current_level:
                level = profile.current_level

        parts.append(f"- Mức năng lực hiện tại: {level}")

        # 2. Điểm yếu — tài liệu có điểm thấp hoặc chưa làm
        if subject_id:
            weak_docs = self.db.query(StudentDocumentEvaluation).filter(
                StudentDocumentEvaluation.user_id == user_id,
                StudentDocumentEvaluation.subject_id == subject_id,
                StudentDocumentEvaluation.latest_score < 50,
            ).order_by(StudentDocumentEvaluation.latest_score.asc()).limit(5).all()

            if weak_docs:
                weak_names = []
                for ev in weak_docs:
                    doc = self.db.query(Document).filter(Document.id == ev.document_id).first()
                    name = self._clean_text(doc.title or doc.filename or f"Tài liệu {ev.document_id}")
                    weak_names.append(f"{name} (điểm {ev.latest_score:.0f})")
                parts.append(f"- Các phần đang yếu: {', '.join(weak_names)}")

            # 3. Misconceptions — câu sai lặp lại gần đây
            from db.models import WrongAnswerRecord
            recent_wrongs = self.db.query(WrongAnswerRecord).filter(
                WrongAnswerRecord.user_id == user_id,
                WrongAnswerRecord.subject_id == subject_id,
            ).order_by(WrongAnswerRecord.created_at.desc()).limit(8).all()

            if recent_wrongs:
                wrong_topics = []
                for w in recent_wrongs[:5]:
                    snippet = self._clean_text(w.question_text or "")[:80]
                    if snippet:
                        wrong_topics.append(f'"{snippet}" → bạn chọn {w.student_choice}, đáp án đúng {w.correct_answer}')
                if wrong_topics:
                    parts.append(
                        f"- Những câu bạn đã làm sai gần đây:\n"
                        + "\n".join(f"  + {t}" for t in wrong_topics)
                    )

            # 4. Strong topics — tài liệu đã vượt qua tốt
            strong_docs = self.db.query(StudentDocumentEvaluation).filter(
                StudentDocumentEvaluation.user_id == user_id,
                StudentDocumentEvaluation.subject_id == subject_id,
                StudentDocumentEvaluation.latest_score >= 75,
            ).order_by(StudentDocumentEvaluation.latest_score.desc()).limit(3).all()

            if strong_docs:
                strong_names = []
                for ev in strong_docs:
                    doc = self.db.query(Document).filter(Document.id == ev.document_id).first()
                    name = self._clean_text(doc.title or doc.filename or f"Tài liệu {ev.document_id}")
                    strong_names.append(f"{name} ({ev.latest_score:.0f}đ)")
                parts.append(f"- Các phần đã nắm vững tốt: {', '.join(strong_names)}")

        # 5. Hành vi học — số lần làm bài
        from db.models import StudentLearningProgress
        progress = self.db.query(StudentLearningProgress).filter_by(user_id=user_id).first()
        if progress:
            total_mins = int(progress.total_study_minutes or 0)
            if total_mins > 0:
                parts.append(f"- Tổng thời gian đã học: {total_mins} phút")

        return "\n".join(parts)

    # ==========================================
    # 2. HÀM GIA SƯ AI CHAT — CÁ NHÂN HÓA
    # ==========================================
    # ================================================================== #
    #  Agentic RAG (tool-calling)                                         #
    #  Thay cho naive RAG: thay vì nhét 22k ký tự tài liệu vào prompt,   #
    #  để LLM tự gọi tool retrieve(query) để lấy đúng đoạn liên quan.     #
    # ================================================================== #
    _RETRIEVE_TOOL = {
        "type": "function",
        "function": {
            "name": "retrieve_document_context",
            "description": (
                "Truy xuất các đoạn tài liệu liên quan đến truy vấn để trả lời câu hỏi của học sinh. "
                "BẮT BUỘC gọi tool này trước khi trả lời; chỉ dùng thông tin trả về để trả lời; "
                "nếu chưa đủ hãy gọi thêm lần nữa với truy vấn cụ thể hơn."
            ),
            "parameters": {
                "type": "object",
                "properties": {
                    "query": {
                        "type": "string",
                        "description": (
                            "Truy vấn tiếng Việt mô tả thông tin cần tìm trong tài liệu "
                            "(ví dụ: 'định nghĩa giao của hai dãy số và ví dụ minh họa')."
                        ),
                    }
                },
                "required": ["query"],
            },
        },
    }

    def _retrieve_chunks(
        self,
        query: str,
        subject: str,
        source_file: str = "",
        allowed_filenames: list = None,
        document_id: Optional[int] = None,
        top_k: int = 6,
    ) -> List[Dict[str, str]]:
        """Semantic search thật theo query của LLM, filter theo tài liệu đang mở."""
        if self.vector_store is None or not (query or "").strip():
            return []

        search_filter = {"subject": {"$eq": subject}}
        if allowed_filenames:
            search_filter = {"$and": [{"subject": {"$eq": subject}}, {"source": {"$in": allowed_filenames}}]}
        if source_file:
            search_filter = {"$and": [{"subject": {"$eq": subject}}, {"source": {"$eq": source_file}}]}
        if document_id:
            doc = self._fetch_document(subject, source_file=source_file, document_id=document_id)
            if doc and doc.filename:
                search_filter = {"$and": [{"subject": {"$eq": subject}}, {"source": {"$eq": doc.filename}}]}

        try:
            docs = self.vector_store.similarity_search(query, k=top_k, filter=search_filter)
            if not docs and source_file:
                docs_all = self.vector_store.similarity_search(query, k=top_k * 3, filter={"subject": {"$eq": subject}})
                docs = [d for d in docs_all if os.path.basename(str((d.metadata or {}).get("source", ""))) == source_file]
        except Exception as exc:
            print(f"⚠️ Agentic retrieve failed: {exc}")
            return []

        chunks: List[Dict[str, str]] = []
        seen = set()
        for doc in docs:
            text = self._clean_text(getattr(doc, "page_content", ""))
            if len(text) < 20:
                continue
            key = text.lower()
            if key in seen:
                continue
            seen.add(key)
            source = (doc.metadata or {}).get("source", source_file or "")
            chunks.append({"source": os.path.basename(str(source)), "content": text})
        return chunks

    def _run_agentic_rag(
        self,
        subject: str,
        user_message: str,
        system_prompt: str,
        history: list,
        source_file: str = "",
        allowed_filenames: list = None,
        document_id: Optional[int] = None,
        max_iterations: int = 3,
    ) -> str:
        """
        Agentic RAG loop: để LLM tự gọi tool retrieve để lấy đúng đoạn tài liệu.
        Trả về câu trả lời cuối cùng (đã bám tài liệu). Fallback nếu vector_store lỗi.
        """
        api_messages: List[Dict] = [{"role": "system", "content": system_prompt}]
        for msg in (history or [])[-8:]:
            role = msg.get("role", "user")
            content = msg.get("content", "")
            if role in ["user", "assistant"] and content:
                api_messages.append({"role": role, "content": content})
        api_messages.append({"role": "user", "content": user_message})

        retrieved_so_far: List[Dict[str, str]] = []

        for iteration in range(max_iterations):
            try:
                completion = self.client.chat.completions.create(
                    messages=api_messages,
                    model=self.model,
                    temperature=0.35,
                    tools=[self._RETRIEVE_TOOL],
                    tool_choice="auto",
                )
            except Exception as exc:
                print(f"⚠️ Agentic LLM call failed at iteration {iteration}: {exc}")
                return ""

            message = completion.choices[0].message
            tool_calls = getattr(message, "tool_calls", None) or []

            if not tool_calls:
                # LLM không cần thêm ngữ cảnh nữa → trả lời cuối.
                return str(getattr(message, "content", "") or "").strip()

            # LLM yêu cầu retrieve: append message assistant (với tool_calls) rồi xử lý từng tool call.
            api_messages.append(message.model_dump())

            for call in tool_calls:
                function = getattr(call, "function", None)
                if not function or getattr(function, "name", "") != "retrieve_document_context":
                    continue
                try:
                    arguments = json.loads(function.arguments or "{}")
                except Exception:
                    arguments = {}
                query = self._clean_text(arguments.get("query", "")) or self._clean_text(user_message)

                chunks = self._retrieve_chunks(
                    query=query,
                    subject=subject,
                    source_file=source_file,
                    allowed_filenames=allowed_filenames,
                    document_id=document_id,
                )
                retrieved_so_far.extend(chunks)

                # Đóng gói context trả về cho LLM (có gắn nguồn để citation).
                if chunks:
                    context_block = "\n\n".join(
                        [f"[Nguồn: {c['source']}]\n{c['content']}" for c in chunks]
                    )
                    tool_content = f"Các đoạn tài liệu liên quan đến truy vấn '{query}':\n\n{context_block}"
                else:
                    tool_content = (
                        f"Không tìm thấy đoạn nào khớp truy vấn '{query}' trong tài liệu. "
                        "Nếu tài liệu không chứa nội dung này, hãy trả lời học sinh là không thấy thông tin trong tài liệu."
                    )

                api_messages.append(
                    {
                        "role": "tool",
                        "tool_call_id": getattr(call, "id", ""),
                        "content": tool_content,
                    }
                )

        # Hết số lần cho phép: ép LLM tổng hợp câu trả lời với context đã có.
        api_messages.append(
            {
                "role": "user",
                "content": "Đã đủ ngữ cảnh. Hãy trả lời câu hỏi ban đầu CHỈ dựa trên các đoạn tài liệu đã truy xuất ở trên. Nếu không có thông tin liên quan, nói rõ là không thấy trong tài liệu.",
            }
        )
        try:
            completion = self.client.chat.completions.create(
                messages=api_messages,
                model=self.model,
                temperature=0.35,
            )
            return str(completion.choices[0].message.content or "").strip()
        except Exception as exc:
            print(f"⚠️ Agentic final answer failed: {exc}")
            return ""

    @staticmethod
    def _tutor_session_key(user_id: Optional[int], subject: str, source_file: str = "") -> Optional[str]:
        """Build a stable memory key per (user, subject, document) so each document keeps its own chat context."""
        if not user_id:
            return None
        return f"tutor|u{user_id}|s{(subject or '').strip().lower()}|d{(source_file or '').strip().lower()}"

    @staticmethod
    def _save_assistant_reply(memory_key: Optional[str], reply: str) -> None:
        """Persist the assistant reply into conversation memory so later turns have context."""
        if not memory_key or not (reply or "").strip():
            return
        try:
            get_conversation_memory().add_message_generic(memory_key, "assistant", reply.strip())
        except Exception as exc:
            print(f"⚠️ Không lưu câu trả lời vào memory: {exc}")

    def chat_with_tutor(self, subject: str, user_message: str, roadmap_context: str, allowed_filenames: list = None, session_topic: str = "", source_file: str = "", history: list = None, document_id: Optional[int] = None, user_id: Optional[int] = None):
        if history is None:
            history = []

        # Vẫn cần tên file (effective_source) để filter retrieval theo đúng tài liệu đang mở.
        document_context = self._get_document_chat_context(
            subject,
            source_file=source_file,
            allowed_filenames=allowed_filenames,
            document_id=document_id,
        )
        document_text = document_context.get("content", "")
        effective_source = document_context.get("source_file", "") or self._clean_text(source_file)

        # --- LƯU NGỮ CẢNH CHAT (ConversationMemory) ---
        # Luôn ưu tiên history do caller truyền (vd Orbit lấy từ DB). Nếu không có,
        # tự load từ ConversationMemory theo (user, subject, document) để giữ context
        # qua nhiều câu hỏi liên tiếp ("giải thích rõ hơn" phải hiểu câu trước đó).
        memory_key = self._tutor_session_key(user_id, subject, effective_source)
        conversation_memory = get_conversation_memory()
        if not history and memory_key:
            try:
                history = conversation_memory.get_history_generic(memory_key)
            except Exception as exc:
                print(f"⚠️ Không load được lịch sử chat: {exc}")
        # Ghi câu hỏi hiện tại vào memory ngay để các lượt sau thấy.
        if memory_key and (user_message or "").strip():
            try:
                conversation_memory.add_message_generic(memory_key, "user", user_message.strip())
            except Exception as exc:
                print(f"⚠️ Không lưu câu hỏi vào memory: {exc}")

        # --- CÁ NHÂN HÓA: Lấy hồ sơ học sinh ---
        student_context = ""
        if user_id:
            try:
                student_context = self._build_student_context(user_id, subject)
            except Exception as exc:
                print(f"⚠️ Không lấy được hồ sơ học sinh: {exc}")

        personalization_block = ""
        if student_context:
            personalization_block = f"""
### HỒ SƠ HỌC SINH (DỰA VÀO DỮ LIỆU THỰC TẾ):
{student_context}

### HƯỚNG DẪN DẠY HỌC CÁ NHÂN HÓA (BẮT BUỘC TUÂN THỦ):
- Điều chỉnh mức độ giải thích theo năng lực học sinh (mục Mức năng lực ở trên).
- Nếu học sinh hỏi về phần đang yếu → giải thích CHẬM, STEP-BY-STEP, cho thêm 1-2 ví dụ minh họa đơn giản.
- Nếu học sinh hỏi về phần đã nắm vững → có thể đi thẳng vào vấn đề, đưa thêm góc nhìn nâng cao hoặc câu hỏi phản biện.
- Nếu câu hỏi liên quan đến misconception đã phát hiện (câu sai gần đây) → CHỦ ĐỘNG kiểm tra xem học sinh có đang hiểu sai không, giải thích sự khác biệt giữa đáp án đúng và sai.
- Giọng văn:
  + Beginner: gần gũi, dùng ví dụ đời thường, tránh thuật ngữ phức tạp.
  + Intermediate: cân bằng giữa lý thuyết và thực hành, khuyến khích tự suy luận.
  + Advanced: tập trung phân tích sâu, thử thách tư duy phản biện, gợi ý mở rộng.
- Không cần nhắc lại profile học sinh trong câu trả lời, chỉ CỨNG DẠY KHÁC NHAU.
"""

        if not self.client:
            return "Tutor Agent hiện chưa kết nối được Groq. Hãy kiểm tra khóa Groq của hệ thống rồi thử lại."

        # --- AGENTIC RAG (ưu tiên): LLM tự gọi tool retrieve để lấy đúng đoạn tài liệu ---
        # Bỏ cách cũ nhét 22k ký tự document_text vào prompt. Chỉ dùng khi vector_store sẵn sàng.
        if self.vector_store is not None:
            agentic_system_prompt = f"""Bạn là Tutor Agent CÁ NHÂN HÓA của môn {subject}.

{personalization_block}
### NHIỆM VỤ CỐT LÕI (BẮT BUỘC):
- TRƯỚC KHI trả lời, BẮT BUỘC gọi tool `retrieve_document_context` với truy vấn mô tả thông tin cần tìm trong tài liệu.
- CHỈ dùng thông tin tool trả về để trả lời. KHÔNG bịa đặt kiến thức ngoài tài liệu.
- Nếu sau khi truy xuất vẫn không thấy thông tin, nói rõ: "Tài liệu hiện tại không đề cập nội dung này".
- Trả lời tiếng Việt, rõ ràng, dễ hiểu, đi thẳng vào yêu cầu.

### NGỮ CẢNH PHIÊN HỌC:
- Buổi học hiện tại: {roadmap_context or session_topic or effective_source or subject}
- Tài liệu đang mở: {effective_source or "Tài liệu hiện tại"}
"""
            started = time.perf_counter()
            trace_prompt = f"USER: {user_message}\n(history {len(history)} lượt, agentic RAG với tool retrieve)"
            log_llm_request("groq", self.model, prompt=trace_prompt, system_prompt=agentic_system_prompt)
            try:
                response_text = self._run_agentic_rag(
                    subject=subject,
                    user_message=user_message,
                    system_prompt=agentic_system_prompt,
                    history=history,
                    source_file=effective_source or source_file,
                    allowed_filenames=allowed_filenames,
                    document_id=document_id,
                )
                duration_ms = (time.perf_counter() - started) * 1000.0
                if response_text:
                    log_llm_response("groq", self.model, response=str(response_text), duration_ms=duration_ms)
                    self._save_assistant_reply(memory_key, response_text)
                    return response_text
                # Agentic trả rỗng (vd LLM lỗi) → fallback sang cách cũ bên dưới.
                print("⚠️ Agentic RAG trả kết quả rỗng, fallback sang naive RAG.")
            except Exception as exc:
                log_llm_error("groq", self.model, error_message=str(exc), duration_ms=0.0)
                print(f"⚠️ Agentic RAG lỗi, fallback sang naive RAG: {exc}")

        # --- FALLBACK (naive RAG, cách cũ): nhét document_text vào prompt ---
        # Dùng khi vector_store bị tắt hoặc agentic loop lỗi, để hệ thống vẫn chạy (degraded).
        if not document_text:
            return "Mình chưa đọc được nội dung của tài liệu đang mở. Hãy mở lại đúng tài liệu này rồi gửi câu hỏi thêm một lần nữa."

        system_prompt = f"""Bạn là Tutor Agent CÁ NHÂN HÓA của môn {subject}.

{personalization_block}
### NHIỆM VỤ CỐT LÕI:
- Trả lời đúng câu hỏi dựa trên tài liệu đang mở — đây là nguồn chính.
- Nếu tài liệu không có đủ thông tin để kết luận, nói rõ là không thấy thông tin đó trong tài liệu hiện tại.
- Không trả lời kiểu chung chung như "không thể hỗ trợ".
- Trả lời bằng tiếng Việt, rõ ràng, dễ hiểu, ưu tiên đi thẳng vào yêu cầu.
- Không cần tự tóm tắt trước nếu người dùng chưa yêu cầu.

### NGỮ CẢNH PHIÊN HỌC:
- Buổi học hiện tại: {roadmap_context or session_topic or effective_source or subject}
- Tài liệu đang mở: {effective_source or "Tài liệu hiện tại"}

[NỘI DUNG TÀI LIỆU ĐANG MỞ]
{document_text}
"""

        api_messages = [{"role": "system", "content": system_prompt}]
        for msg in history[-8:]:
            role = msg.get("role", "user")
            content = msg.get("content", "")
            if role in ["user", "assistant"] and content:
                api_messages.append({"role": role, "content": content})
        api_messages.append({"role": "user", "content": user_message})

        try:
            trace_prompt = "\n\n".join([
                f"{msg.get('role', 'user').upper()}: {msg.get('content', '')}"
                for msg in api_messages[1:]
                if str(msg.get('content', '')).strip()
            ])
            started = time.perf_counter()
            log_llm_request("groq", self.model, prompt=trace_prompt, system_prompt=system_prompt)
            chat_completion = self.client.chat.completions.create(
                messages=api_messages,
                model=self.model,
                temperature=0.35,
            )
            response_text = str(chat_completion.choices[0].message.content or "").strip()
            duration_ms = (time.perf_counter() - started) * 1000.0
            log_llm_response("groq", self.model, response=str(response_text), duration_ms=duration_ms)
            self._save_assistant_reply(memory_key, response_text)
            return response_text or "Mình chưa nhận được câu trả lời từ Groq cho tài liệu này. Bạn hãy gửi lại câu hỏi ngắn hơn."
        except Exception as e:
            log_llm_error("groq", self.model, error_message=str(e), duration_ms=0.0)
            error_text = str(e).lower()
            if "invalid_api_key" in error_text or "401" in error_text:
                return "Tutor Agent chưa thể gọi Groq vì khóa API không hợp lệ hoặc đã hết hiệu lực."
            if "rate limit" in error_text or "429" in error_text:
                return "Tutor Agent đang chạm giới hạn tạm thời từ Groq. Bạn thử gửi lại sau ít phút."
            return "Tutor Agent đang lỗi kết nối tới Groq. Bạn hãy thử gửi lại câu hỏi ngắn hơn."

    # ==========================================
    # 3. HÀM SINH CÂU HỎI TRẮC NGHIỆM
    # ==========================================
    def _build_session_quiz_fallback(self, subject: str, session_topic: str, level: str, context_docs: str, num_questions: int):
        raw_sentences = re.split(r"(?<=[\.!\?;:])\s+", context_docs or "")
        sentences = []
        for s in raw_sentences:
            clean = re.sub(r"\s+", " ", s).strip()
            if 55 <= len(clean) <= 260:
                sentences.append(clean)

        if not sentences:
            sentences = [f"Buổi học {session_topic} của môn {subject} tập trung vào kiến thức cốt lõi và cách vận dụng trong bài tập."]

        token_bank = re.findall(r"[A-Za-zÀ-ỹ0-9_]{4,}", context_docs or "")
        token_bank = [t for t in token_bank if not t.isdigit()]
        token_bank = list(dict.fromkeys(token_bank))
        numeric_bank = list(dict.fromkeys(re.findall(r"\d+(?:[\.,]\d+)?", context_docs or "")))

        antonym_pairs = [
            ("tăng", "giảm"),
            ("đúng", "sai"),
            ("lớn hơn", "nhỏ hơn"),
            ("trước", "sau"),
            ("hội tụ", "phân kỳ"),
            ("đồng biến", "nghịch biến"),
            ("liên tục", "gián đoạn"),
        ]

        def mutate_number(text: str) -> str:
            nums = re.findall(r"\d+(?:[\.,]\d+)?", text)
            if not nums:
                return ""
            old = nums[0]
            replacement = None
            cands = [n for n in numeric_bank if n != old]
            if cands:
                replacement = random.choice(cands)
            if not replacement:
                try:
                    replacement = str(int(float(old.replace(",", "."))) + random.choice([-2, -1, 1, 2]))
                except Exception:
                    replacement = old + "1"
            return re.sub(re.escape(old), replacement, text, count=1)

        def mutate_term(text: str) -> str:
            words = re.findall(r"[A-Za-zÀ-ỹ0-9_]{5,}", text)
            if not words:
                return ""
            pick = words[min(1, len(words) - 1)]
            cands = [t for t in token_bank if t.lower() != pick.lower() and abs(len(t) - len(pick)) <= 6]
            if not cands:
                return ""
            return re.sub(re.escape(pick), random.choice(cands), text, count=1)

        def mutate_relation(text: str) -> str:
            out = text
            for a, b in antonym_pairs:
                if re.search(re.escape(a), out, flags=re.IGNORECASE):
                    return re.sub(re.escape(a), b, out, count=1, flags=re.IGNORECASE)
                if re.search(re.escape(b), out, flags=re.IGNORECASE):
                    return re.sub(re.escape(b), a, out, count=1, flags=re.IGNORECASE)
            return ""

        templates = [
            "Theo học liệu buổi '{topic}', phát biểu nào đúng nhất về ý: \"{anchor}\"?",
            "Dựa trên tài liệu của buổi '{topic}', phương án nào phản ánh chính xác nội dung: \"{anchor}\"?",
            "Trong ngữ cảnh môn {subject}, nhận định nào đúng với nội dung đã học: \"{anchor}\"?",
        ]

        results = []
        for i in range(max(1, num_questions)):
            fact = sentences[i % len(sentences)]
            correct = fact[:210]
            near = sentences[(i + 1) % len(sentences)][:210]

            d1 = mutate_number(correct)
            d2 = mutate_term(correct)
            d3 = mutate_relation(correct)

            distractors = []
            for d in [d1, d2, d3, near]:
                d_clean = (d or "").strip()[:210]
                if d_clean and d_clean.lower() != correct.lower() and d_clean.lower() not in [x.lower() for x in distractors]:
                    distractors.append(d_clean)
                if len(distractors) == 3:
                    break

            while len(distractors) < 3:
                distractors.append((correct[:120] + " nhưng điều kiện áp dụng không khớp với tài liệu.")[:210])

            options_raw = [correct] + distractors[:3]
            order = [0, 1, 2, 3]
            random.shuffle(order)
            labels = ["A", "B", "C", "D"]
            final_options = []
            correct_label = "A"
            for pos, idx in enumerate(order):
                final_options.append(f"{labels[pos]}. {options_raw[idx]}")
                if idx == 0:
                    correct_label = labels[pos]

            q_template = templates[i % len(templates)]
            question_text = q_template.format(topic=session_topic, anchor=fact[:110], subject=subject)

            results.append(
                {
                    "id": i + 1,
                    "content": question_text,
                    "options": final_options,
                    "correct_label": correct_label,
                    "explanation": f"Câu fallback theo mức {str(level).upper()}, dựa trên nội dung tài liệu đã học; phương án nhiễu thay đổi thuật ngữ/số liệu/quan hệ.",
                }
            )

        return results

    def generate_session_quiz(self, subject: str, session_topic: str, level: str, allowed_filenames: list = None):
        context_docs = ""
        search_filter = {"subject": {"$eq": subject}}
        if allowed_filenames:
            search_filter = {
                "$and": [
                    {"subject": {"$eq": subject}},
                    {"source": {"$in": allowed_filenames}}
                ]
            }

        is_final_exam = "CUỐI KHÓA" in session_topic.upper() or "TỔNG HỢP" in session_topic.upper()

        # PHÂN LUỒNG TÌM KIẾM ĐỂ TRÁNH LẠC ĐỀ
        if is_final_exam:
            # Bốc diện rộng toàn bộ các chương
            search_query = f"Toàn bộ kiến thức trọng tâm, các khái niệm, bài tập, ứng dụng và tổng kết của môn {subject}"
            topic_instruction = f"- Chủ đề: TỔNG ÔN CUỐI KHÓA (Lệnh: Bốc ngẫu nhiên kiến thức rải rác ở TẤT CẢ CÁC CHƯƠNG để kiểm tra toàn diện)."
            num_questions = 20
            k_val = 30 # Lấy tận 30 mảnh kiến thức khác nhau
        else:
            # Focus cực mạnh vào đúng 1 keyword
            search_query = f"Kiến thức chuyên sâu, ví dụ thực tiễn, đoạn code, bài tập của chủ đề: {session_topic} trong môn {subject}"
            topic_instruction = f"- Chủ đề: {session_topic} (Lệnh Sống Còn: TUYỆT ĐỐI CHỈ HỎI XOAY QUANH CHỦ ĐỀ NÀY, không hỏi lan man chương khác)."
            num_questions = 10
            k_val = 20

        if self.vector_store is not None:
            try:
                docs = self.vector_store.similarity_search(search_query, k=k_val, filter=search_filter)
                if not docs and allowed_filenames:
                    docs_all = self.vector_store.similarity_search(search_query, k=max(k_val * 2, 40))
                    docs = [d for d in docs_all if os.path.basename(d.metadata.get("source", "")) in allowed_filenames]
                # Ép dung lượng nạp lên mức 15.000 ký tự để AI thông minh nhất có thể
                context_docs = "\n\n".join([doc.page_content for doc in docs])[:15000]
            except Exception as e:
                print(f"❌ Lỗi RAG lấy tài liệu: {e}")
                context_docs = "Dữ liệu kiến thức đang được cập nhật."
        else:
            context_docs = "Dữ liệu kiến thức đang được cập nhật."

        prompt = f"""
        BẠN LÀ CHUYÊN GIA KHẢO THÍ ĐẠI HỌC CỰC KỲ KHẮT KHE. NHIỆM VỤ: Soạn ĐÚNG {num_questions} câu hỏi trắc nghiệm trình độ {level.upper()}.
        - Môn học: {subject}
        {topic_instruction}

        [TÀI LIỆU CỐT LÕI (BẮT BUỘC BÁM SÁT)]:
        {context_docs}

        [TIÊU CHUẨN CHẤT LƯỢNG ĐỀ THI (CHỐNG HỌC VẸT)]:
        1. KHÔNG hỏi lý thuyết suông (Vd: "Định nghĩa là gì?", "Cấu trúc gồm mấy phần?"). Sinh viên Đại học cần tư duy!
        2. BẮT BUỘC sử dụng Tình huống thực tiễn (Case study), Đoạn code/Công thức, hoặc Bài toán logic để sinh viên phân tích và giải quyết.
        3. Phân cấp độ:
           - BEGINNER: Nhận biết khái niệm thông qua ví dụ thực tế.
           - INTERMEDIATE: Phân tích đúng sai, dự đoán kết quả, tìm lỗi sai.
           - ADVANCED: Đánh giá, tối ưu hóa hệ thống, giải quyết tình huống hóc búa.
        4. TÍNH KHÁCH QUAN: 1 đáp án ĐÚNG CHÍNH XÁC, 3 đáp án SAI NHƯNG CÓ VẺ ĐÚNG (bẫy nhiễu).

        [YÊU CẦU ĐẦU RA JSON]:
        {{
            "questions": [
                {{
                    "id": 1,
                    "content": "Tình huống / Đoạn mã / Câu hỏi tư duy...",
                    "options": ["A. Đáp án 1", "B. Đáp án 2", "C. Đáp án 3", "D. Đáp án 4"],
                    "correct_label": "A", 
                    "explanation": "Giải thích chi tiết vì sao đúng và phân tích bẫy sai..."
                }}
            ]
        }}
        """

        if not self.client:
            return self._build_session_quiz_fallback(subject, session_topic, level, context_docs, num_questions)

        try:
            chat_completion = self.client.chat.completions.create(
                messages=[
                    {"role": "system", "content": "You are a top-tier examiner. You output strict and valid JSON containing the exact number of requested questions."},
                    {"role": "user", "content": prompt}
                ],
                model=self.model,
                max_tokens=6000, # Bơm token để viết đề dài
                temperature=0.3, # Tăng nhẹ độ sáng tạo để bớt rập khuôn
                response_format={"type": "json_object"}
            )
            
            result = json.loads(chat_completion.choices[0].message.content)
            questions = result.get("questions", [])
            if not questions:
                return self._build_session_quiz_fallback(subject, session_topic, level, context_docs, num_questions)
            return questions
        except Exception as e:
            print(f"❌ Lỗi sinh đề thi: {e}")
            return self._build_session_quiz_fallback(subject, session_topic, level, context_docs, num_questions)
