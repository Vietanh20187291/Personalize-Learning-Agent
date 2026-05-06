import os
import json
import re
import logging
import time
from pathlib import Path
from fastapi import APIRouter, Depends, HTTPException
from fastapi.responses import FileResponse, StreamingResponse 
from sqlalchemy.orm import Session
from pydantic import BaseModel
from db.database import get_db
from db import models
from typing import List, Optional
from datetime import timedelta
from langchain_community.document_loaders import PyPDFLoader, Docx2txtLoader, TextLoader
from pptx import Presentation

# Import hàm lấy vector store
from rag.vector_store import get_vector_store 
from agents.assessment_agent import AssessmentAgent
from config import settings

router = APIRouter()
logger = logging.getLogger("app.document")
_PREVIEW_CACHE: dict[str, tuple[float, dict]] = {}


def _sanitize_question_text(text: Optional[str]) -> str:
    cleaned = re.sub(r"\s+", " ", str(text or "")).strip()
    cleaned = re.sub(r"\s+#\d+\b", "", cleaned)
    cleaned = re.sub(r"\s*\(\s*mức\s*cơ\s*bản\s*,\s*câu\s*\d+\s*\)", "", cleaned, flags=re.IGNORECASE)
    cleaned = re.sub(
        r"\s*(?:Trá»ng tÃ¢m|Trong ngá»¯ cáº£nh|Khi xÃ©t|Trong pháº¡m vi)\s*:\s*[^.?!;]+[.?!;:]*\s*$",
        "",
        cleaned,
        flags=re.IGNORECASE,
    )
    return re.sub(r"\s+", " ", cleaned).strip()


def _preview_cache_key(doc_id: int, filename: str) -> str:
    return f"{doc_id}:{(filename or '').strip()}"


def _get_cached_preview(cache_key: str) -> Optional[dict]:
    ttl_seconds = max(60, int(getattr(settings, "PREVIEW_CACHE_TTL_SECONDS", 900)))
    cached = _PREVIEW_CACHE.get(cache_key)
    if not cached:
        return None

    cached_at, payload = cached
    if time.time() - cached_at > ttl_seconds:
        _PREVIEW_CACHE.pop(cache_key, None)
        return None
    return payload


def _set_cached_preview(cache_key: str, payload: dict) -> None:
    _PREVIEW_CACHE[cache_key] = (time.time(), payload)


def _normalize_question_options(raw_options):
    """Ensure QuestionBank options are always returned as a list of strings."""
    if isinstance(raw_options, list):
        return [_sanitize_question_text(str(opt)) for opt in raw_options]

    if isinstance(raw_options, str):
        try:
            parsed = json.loads(raw_options)
            if isinstance(parsed, list):
                return [_sanitize_question_text(str(opt)) for opt in parsed]
        except Exception:
            pass

    return []


def _cleanup_question_bank_rows(rows: List[models.QuestionBank], db: Session) -> None:
    changed = False
    for row in rows:
        existing_options = _normalize_question_options(getattr(row, "options", []))
        cleaned_content = _sanitize_question_text(getattr(row, "content", ""))
        cleaned_explanation = _sanitize_question_text(re.sub(r"^\[bloom:[^\]]+\]\s*", "", getattr(row, "explanation", "") or "", flags=re.IGNORECASE))
        cleaned_options = [_sanitize_question_text(opt) for opt in existing_options]
        if cleaned_content and cleaned_content != (row.content or ""):
            row.content = cleaned_content
            changed = True
        original_explanation = re.sub(r"^\[bloom:[^\]]+\]\s*", "", getattr(row, "explanation", "") or "", flags=re.IGNORECASE).strip()
        if cleaned_explanation != original_explanation:
            bloom_prefix_match = re.match(r"^(\[bloom:[^\]]+\]\s*)", getattr(row, "explanation", "") or "", flags=re.IGNORECASE)
            bloom_prefix = bloom_prefix_match.group(1) if bloom_prefix_match else ""
            row.explanation = f"{bloom_prefix}{cleaned_explanation}".strip()
            changed = True
        if cleaned_options and cleaned_options != existing_options:
            row.options = cleaned_options
            changed = True
    if changed:
        db.commit()


def _resolve_document_file_path(filename: str) -> Path:
    """Resolve uploaded file path independent of current working directory."""
    safe_name = (filename or "").strip()
    if not safe_name:
        raise HTTPException(status_code=404, detail="Tên file tài liệu không hợp lệ")

    project_root = Path(__file__).resolve().parents[2]
    candidates = [
        project_root / "temp_uploads" / safe_name,
        project_root / "backend" / "temp_uploads" / safe_name,
    ]

    for candidate in candidates:
        if candidate.exists() and candidate.is_file():
            return candidate

    raise HTTPException(status_code=404, detail="File vật lý không tồn tại trên server")


def _can_student_access_document(db: Session, doc: models.Document, user_id: Optional[int]) -> bool:
    publication = db.query(models.DocumentPublication).filter(models.DocumentPublication.doc_id == doc.id).first()
    if publication and publication.is_visible_to_students:
        return True

    if not user_id:
        return False

    user = db.query(models.User).filter(models.User.id == user_id).first()
    if not user or getattr(user, "role", None) != "student":
        return False

    enrolled_class_ids = {classroom.id for classroom in (getattr(user, "enrolled_classes", []) or [])}
    return bool(doc.class_id and doc.class_id in enrolled_class_ids)


def _extract_preview_segments(file_path: str, filename: str):
    """Trích xuất nội dung text từ các file để hiển thị preview.
    Format tốt hơn để presentation rõ ràng.
    """
    ext = os.path.splitext(filename)[1].lower()
    
    try:
        if ext == ".pptx":
            slides = []
            try:
                prs = Presentation(file_path)
                for i, slide in enumerate(prs.slides, start=1):
                    slide_lines = []
                    
                    for shape in slide.shapes:
                        try:
                            if hasattr(shape, "text_frame"):
                                for paragraph in shape.text_frame.paragraphs:
                                    text = " ".join(paragraph.text.split()).strip()
                                    if text:
                                        slide_lines.append(text)
                            elif hasattr(shape, "text") and isinstance(shape.text, str):
                                text = " ".join(shape.text.split()).strip()
                                if text:
                                    slide_lines.append(text)
                        except Exception as shape_err:
                            continue
                    
                    if slide_lines:
                        # Format better: join with bullet points
                        formatted_content = "\n".join([f"• {line}" for line in slide_lines])
                        slides.append({
                            "title": f"📊 Slide {i}",
                            "content": formatted_content
                        })
                
                if slides:
                    return {"type": "pptx", "segments": slides}
                else:
                    return {"type": "pptx", "segments": []}
                    
            except Exception as pptx_err:
                print(f"❌ Lỗi parse PPTX '{filename}': {pptx_err}")
                return {"type": "pptx", "segments": []}

        if ext == ".pdf":
            pages = []
            try:
                loader = PyPDFLoader(file_path)
                docs = loader.load()
                
                for idx, doc in enumerate(docs, start=1):
                    text = " ".join((doc.page_content or "").split()).strip()
                    if text:
                        pages.append({
                            "title": f"📄 Trang {idx}",
                            "content": text[:500] + ("..." if len(text) > 500 else "")
                        })
                
                if pages:
                    return {"type": "pdf", "segments": pages}
                else:
                    return {"type": "pdf", "segments": []}
                    
            except Exception as pdf_err:
                print(f"❌ Lỗi parse PDF '{filename}': {pdf_err}")
                return {"type": "pdf", "segments": []}

        if ext == ".docx":
            try:
                loader = Docx2txtLoader(file_path)
                docs = loader.load()
                content = "\n".join([(d.page_content or "").strip() for d in docs if (d.page_content or "").strip()])
                
                blocks = [b.strip() for b in content.split("\n") if b.strip()]
                
                if blocks:
                    segments = [{"title": f"📋 Đoạn {i+1}", "content": b} for i, b in enumerate(blocks[:150])]
                    return {"type": "docx", "segments": segments}
                else:
                    return {"type": "docx", "segments": []}
                    
            except Exception as docx_err:
                print(f"❌ Lỗi parse DOCX '{filename}': {docx_err}")
                return {"type": "docx", "segments": []}

        if ext == ".txt":
            try:
                loader = TextLoader(file_path, encoding="utf-8")
                docs = loader.load()
                content = "\n".join([(d.page_content or "").strip() for d in docs if (d.page_content or "").strip()])
                
                blocks = [b.strip() for b in content.split("\n") if b.strip()]
                
                if blocks:
                    segments = [{"title": f"📝 Dòng {i+1}", "content": b} for i, b in enumerate(blocks[:200])]
                    return {"type": "txt", "segments": segments}
                else:
                    return {"type": "txt", "segments": []}
                    
            except Exception as txt_err:
                print(f"⚠️ Lỗi parse TXT (UTF-8) '{filename}': {txt_err}")
                try:
                    loader = TextLoader(file_path, encoding="cp1252")
                    docs = loader.load()
                    content = "\n".join([(d.page_content or "").strip() for d in docs if (d.page_content or "").strip()])
                    blocks = [b.strip() for b in content.split("\n") if b.strip()]
                    
                    if blocks:
                        segments = [{"title": f"📝 Dòng {i+1}", "content": b} for i, b in enumerate(blocks[:200])]
                        return {"type": "txt", "segments": segments}
                except:
                    pass
                return {"type": "txt", "segments": []}

        return {"type": "unknown", "segments": []}
        
    except Exception as general_err:
        print(f"❌ Lỗi chung khi parse '{filename}': {general_err}")
        return {"type": "unknown", "segments": []}


def _extract_preview_from_vector(source_filename: str):
    try:
        vector_store = get_vector_store()
        collection = vector_store._collection
        data = collection.get(include=["documents", "metadatas"])
        docs = data.get("documents", []) or []
        metas = data.get("metadatas", []) or []

        matched = []
        for i, meta in enumerate(metas):
            source = ""
            if meta and isinstance(meta, dict):
                source = str(meta.get("source", ""))
            if os.path.basename(source) == source_filename:
                txt = str(docs[i] or "").strip()
                if txt:
                    matched.append(txt)

        segments = [
            {"title": f"Nội dung {idx + 1}", "content": text}
            for idx, text in enumerate(matched[:80])
        ]
        return segments
    except Exception:
        return []

class QuestionBankCreateRequest(BaseModel):
    content: str
    options: List[str]
    correct_answer: str
    explanation: Optional[str] = None
    difficulty: Optional[str] = None

class QuestionBankUpdateRequest(BaseModel):
    content: Optional[str] = None
    options: Optional[List[str]] = None
    correct_answer: Optional[str] = None
    explanation: Optional[str] = None
    difficulty: Optional[str] = None


def _resolve_document_file_path(filename: str) -> Path:
    """Resolve uploaded file path independent of current working directory."""
    safe_name = (filename or "").strip()
    if not safe_name:
        raise HTTPException(status_code=404, detail="Tên file tài liệu không hợp lệ")

    project_root = Path(__file__).resolve().parents[2]
    candidates = [
        project_root / "temp_uploads" / safe_name,
        project_root / "backend" / "temp_uploads" / safe_name,
    ]

    for candidate in candidates:
        if candidate.exists() and candidate.is_file():
            return candidate

    raise HTTPException(status_code=404, detail="File vật lý không tồn tại trên server")


@router.get("/class-documents/{class_id}")
def get_documents(
    class_id: int, 
    subject: Optional[str] = None, 
    db: Session = Depends(get_db)
):
    """Lấy danh sách tài liệu của một lớp với múi giờ VN"""
    query = db.query(models.Document).filter(models.Document.class_id == class_id)
    
    if subject and subject != "all":
        query = query.filter(models.Document.subject == subject)
        
    docs = query.order_by(models.Document.upload_time.desc()).all()
    
    results = []
    for doc in docs:
        # FIX GIỜ VIỆT NAM (UTC+7)
        time_str = (doc.upload_time + timedelta(hours=7)).strftime("%H:%M - %d/%m/%Y") if doc.upload_time else "Chưa xác định"
        
        results.append({
            "id": doc.id,
            "title": doc.filename, 
            "subject": doc.subject,
            "file_path": f"temp_uploads/{doc.filename}", 
            "created_at": time_str 
        })
        
    return results

# ==========================================
# 2. API CHO GIÁO VIÊN: XÓA TÀI LIỆU
# ==========================================
@router.delete("/delete/{doc_id}")
def delete_document(doc_id: int, db: Session = Depends(get_db)):
    """Xóa tài liệu hoàn toàn khỏi 3 nơi: File vật lý, ChromaDB và SQLite"""
    doc = db.query(models.Document).filter(models.Document.id == doc_id).first()
    if not doc:
        raise HTTPException(status_code=404, detail="Không tìm thấy tài liệu")
    
    filename = doc.filename

    # --- 1. XÓA FILE VẬT LÝ TRONG THƯ MỤC TEMP_UPLOADS ---
    if filename:
        file_path = os.path.join("temp_uploads", filename)
        if os.path.exists(file_path):
            try:
                os.remove(file_path)
                print(f"✅ Đã dọn dẹp file vật lý: {file_path}")
            except Exception as e:
                print(f"❌ Lỗi khi xóa file vật lý: {e}")

    # --- 2. XÓA TRI THỨC TRONG CHROMADB ---
    try:
        vector_store = get_vector_store()

        # Lấy trực tiếp đối tượng Collection gốc của ChromaDB thay vì dùng vỏ bọc LangChain
        collection = vector_store._collection

        # Ép Collection phải trả về TOÀN BỘ dữ liệu metadatas (không bị limit)
        all_data = collection.get(include=["metadatas"])
        ids = all_data.get("ids", [])
        metadatas = all_data.get("metadatas", [])

        ids_to_delete = []
        for i, meta in enumerate(metadatas):
            if meta and "source" in meta:
                source_path = str(meta["source"])
                # Miễn là tên file có xuất hiện trong đường dẫn thì gom vào danh sách tử hình
                if filename in source_path:
                    ids_to_delete.append(ids[i])

        if ids_to_delete:
            # Lệnh xóa ép buộc từ Core ChromaDB
            collection.delete(ids=ids_to_delete)
            print(f"✅ Đã dọn dẹp sạch {len(ids_to_delete)} đoạn băm của {filename} trong ChromaDB")
        else:
            print(f"⚠️ Không tìm thấy tri thức của {filename} trong ChromaDB")

    except Exception as e:
        print(f"⚠️ Bỏ qua dọn ChromaDB cho '{filename}' để giữ backend ổn định: {e}")

    # --- 3. XÓA DỮ LIỆU TRONG SQLITE ---
    try:
        # Xóa các Chunks liên quan trong bảng chunks (nếu có)
        db.query(models.Chunk).filter(models.Chunk.source_file == filename).delete()
        
        # Xóa bản ghi trong bảng documents
        db.delete(doc)
        db.commit()
        return {"message": "Đã xóa tài liệu và tri thức AI thành công"}
    except Exception as e:
        db.rollback()
        raise HTTPException(status_code=500, detail=f"Lỗi database: {str(e)}")


@router.post("/generate-question-bank/{doc_id}")
def generate_question_bank_for_document(
    doc_id: int,
    target_count: int = 100,
    db: Session = Depends(get_db),
):
    """Tạo lại bộ câu hỏi cho 1 tài liệu: xóa câu hỏi cũ theo tài liệu rồi sinh mới."""
    started = time.perf_counter()
    doc = db.query(models.Document).filter(models.Document.id == doc_id).first()
    if not doc:
        raise HTTPException(status_code=404, detail="Không tìm thấy tài liệu")

    source_file = (doc.filename or "").strip()
    if not source_file:
        raise HTTPException(status_code=400, detail="Tài liệu không có filename hợp lệ")

    subject_id = getattr(doc, "subject_id", None)
    subject_name = (doc.subject or "").strip()
    if subject_id:
        subject_obj = db.query(models.Subject).filter(models.Subject.id == subject_id).first()
        resolved_name = (getattr(subject_obj, "name", "") or "").strip()
        if resolved_name:
            subject_name = resolved_name

    if not subject_name:
        raise HTTPException(status_code=400, detail="Không xác định được môn học của tài liệu")

    safe_target = max(1, min(int(target_count), 60))
    logger.info(
        "generate_question_bank start doc_id=%s subject_id=%s subject=%s source_file=%s target_count=%s",
        doc_id,
        subject_id,
        subject_name,
        source_file,
        safe_target,
    )

    # Xóa toàn bộ câu hỏi cũ đã sinh cho tài liệu này.
    qb_delete_query = db.query(models.QuestionBank).filter(
        models.QuestionBank.source_file == source_file,
    )
    if subject_id:
        qb_delete_query = qb_delete_query.filter(models.QuestionBank.subject_id == subject_id)
    else:
        qb_delete_query = qb_delete_query.filter(models.QuestionBank.subject == subject_name)

    deleted_count = qb_delete_query.delete(synchronize_session=False)
    db.commit()

    agent = AssessmentAgent(db)
    generated = agent.pre_generate_questions_for_document(
        subject=subject_name,
        source_file=source_file,
        count=safe_target,
        force_refresh=True,
        replace_existing=True,
    )

    qb_count_query = db.query(models.QuestionBank).filter(
        models.QuestionBank.source_file == source_file,
    )
    if subject_id:
        qb_count_query = qb_count_query.filter(models.QuestionBank.subject_id == subject_id)
    else:
        qb_count_query = qb_count_query.filter(models.QuestionBank.subject == subject_name)

    final_count = qb_count_query.count()
    logger.info(
        "generate_question_bank done doc_id=%s source_file=%s deleted_count=%s generated_count=%s final_count=%s duration_ms=%s",
        doc_id,
        source_file,
        deleted_count,
        len(generated),
        final_count,
        round((time.perf_counter() - started) * 1000, 2),
    )

    return {
        "doc_id": doc_id,
        "subject": subject_name,
        "source_file": source_file,
        "deleted_count": int(deleted_count),
        "generated_count": int(len(generated)),
        "final_count": int(final_count),
        "target_count": safe_target,
    }


@router.get("/question-bank/{doc_id}")
def get_question_bank_for_document(doc_id: int, db: Session = Depends(get_db)):
    doc = db.query(models.Document).filter(models.Document.id == doc_id).first()
    if not doc:
        raise HTTPException(status_code=404, detail="Không tìm thấy tài liệu")

    source_file = (doc.filename or "").strip()
    if not source_file:
        raise HTTPException(status_code=400, detail="Tài liệu không có filename hợp lệ")

    subject_id = getattr(doc, "subject_id", None)
    subject_name = (doc.subject or "").strip()
    if subject_id:
        subject_obj = db.query(models.Subject).filter(models.Subject.id == subject_id).first()
        resolved_name = (getattr(subject_obj, "name", "") or "").strip()
        if resolved_name:
            subject_name = resolved_name

    if not subject_name:
        raise HTTPException(status_code=400, detail="Không xác định được môn học của tài liệu")

    qb_query = db.query(models.QuestionBank).filter(
        models.QuestionBank.source_file == source_file,
    )
    if subject_id:
        qb_query = qb_query.filter(models.QuestionBank.subject_id == subject_id)
    else:
        qb_query = qb_query.filter(models.QuestionBank.subject == subject_name)

    questions = qb_query.order_by(models.QuestionBank.id.desc()).all()
    _cleanup_question_bank_rows(questions, db)

    return [
        {
            "id": q.id,
            "content": q.content,
            "options": _normalize_question_options(q.options),
            "correct_answer": q.correct_answer,
            "explanation": q.explanation,
            "difficulty": q.difficulty,
            "subject": q.subject,
            "source_file": q.source_file,
        }
        for q in questions
    ]


@router.post("/question-bank/{doc_id}")
def add_question_bank_for_document(
    doc_id: int,
    req: QuestionBankCreateRequest,
    db: Session = Depends(get_db),
):
    doc = db.query(models.Document).filter(models.Document.id == doc_id).first()
    if not doc:
        raise HTTPException(status_code=404, detail="Không tìm thấy tài liệu")

    source_file = (doc.filename or "").strip()
    if not source_file:
        raise HTTPException(status_code=400, detail="Tài liệu không có filename hợp lệ")

    subject_name = (doc.subject or "").strip()
    subject_id = getattr(doc, "subject_id", None)
    if not subject_name and subject_id:
        subject_obj = db.query(models.Subject).filter(models.Subject.id == subject_id).first()
        subject_name = (getattr(subject_obj, "name", "") or "").strip()

    if not subject_name or not subject_id:
        raise HTTPException(status_code=400, detail="Không xác định được môn học của tài liệu")

    validation_options = [opt.strip() for opt in req.options if isinstance(opt, str) and opt.strip()]
    if len(validation_options) < 2:
        raise HTTPException(status_code=400, detail="Vui lòng cung cấp ít nhất 2 lựa chọn hợp lệ")

    new_q = models.QuestionBank(
        subject_id=subject_id,
        subject=subject_name,
        source_file=source_file,
        content=_sanitize_question_text(req.content),
        options=[_sanitize_question_text(opt) for opt in validation_options],
        correct_answer=req.correct_answer.strip(),
        explanation=_sanitize_question_text(req.explanation or ""),
        difficulty=(req.difficulty or "").strip() or "Medium",
    )
    db.add(new_q)
    db.commit()
    db.refresh(new_q)

    return {
        "id": new_q.id,
        "content": new_q.content,
        "options": new_q.options,
        "correct_answer": new_q.correct_answer,
        "explanation": new_q.explanation,
        "difficulty": new_q.difficulty,
        "subject": new_q.subject,
        "source_file": new_q.source_file,
    }


@router.put("/question-bank/{question_id}")
def update_question_bank_entry(
    question_id: int,
    req: QuestionBankUpdateRequest,
    db: Session = Depends(get_db),
):
    q = db.query(models.QuestionBank).filter(models.QuestionBank.id == question_id).first()
    if not q:
        raise HTTPException(status_code=404, detail="Không tìm thấy câu hỏi")

    if req.content is not None:
        q.content = _sanitize_question_text(req.content)
    if req.options is not None:
        validation_options = [opt.strip() for opt in req.options if isinstance(opt, str) and opt.strip()]
        if len(validation_options) < 2:
            raise HTTPException(status_code=400, detail="Vui lòng cung cấp ít nhất 2 lựa chọn hợp lệ")
        q.options = [_sanitize_question_text(opt) for opt in validation_options]
    if req.correct_answer is not None:
        q.correct_answer = req.correct_answer.strip()
    if req.explanation is not None:
        q.explanation = _sanitize_question_text(req.explanation)
    if req.difficulty is not None:
        q.difficulty = req.difficulty.strip()

    db.commit()
    db.refresh(q)

    return {
        "id": q.id,
        "content": q.content,
        "options": q.options,
        "correct_answer": q.correct_answer,
        "explanation": q.explanation,
        "difficulty": q.difficulty,
        "subject": q.subject,
        "source_file": q.source_file,
    }


@router.delete("/question-bank/{question_id}")
def delete_question_bank_entry(question_id: int, db: Session = Depends(get_db)):
    q = db.query(models.QuestionBank).filter(models.QuestionBank.id == question_id).first()
    if not q:
        raise HTTPException(status_code=404, detail="Không tìm thấy câu hỏi")

    db.delete(q)
    db.commit()
    return {"message": "Câu hỏi đã được xóa", "id": question_id}


@router.post("/generate-question-bank/{doc_id}/append")
def append_question_bank_for_document(
    doc_id: int,
    count: int = 10,
    db: Session = Depends(get_db),
):
    started = time.perf_counter()
    doc = db.query(models.Document).filter(models.Document.id == doc_id).first()
    if not doc:
        raise HTTPException(status_code=404, detail="Không tìm thấy tài liệu")

    source_file = (doc.filename or "").strip()
    if not source_file:
        raise HTTPException(status_code=400, detail="Tài liệu không có filename hợp lệ")

    subject_id = getattr(doc, "subject_id", None)
    subject_name = (doc.subject or "").strip()
    if subject_id:
        subject_obj = db.query(models.Subject).filter(models.Subject.id == subject_id).first()
        resolved_name = (getattr(subject_obj, "name", "") or "").strip()
        if resolved_name:
            subject_name = resolved_name

    if not subject_name:
        raise HTTPException(status_code=400, detail="Không xác định được môn học của tài liệu")

    safe_count = max(1, min(int(count), 50))
    logger.info(
        "append_question_bank start doc_id=%s subject_id=%s subject=%s source_file=%s requested_count=%s",
        doc_id,
        subject_id,
        subject_name,
        source_file,
        safe_count,
    )

    existing_query = db.query(models.QuestionBank).filter(
        models.QuestionBank.source_file == source_file,
    )
    if subject_id:
        existing_query = existing_query.filter(models.QuestionBank.subject_id == subject_id)
    else:
        existing_query = existing_query.filter(models.QuestionBank.subject == subject_name)

    existing_rows = existing_query.order_by(models.QuestionBank.id.desc()).all()
    _cleanup_question_bank_rows(existing_rows, db)
    existing_count = len(existing_rows)

    agent = AssessmentAgent(db)
    generated = agent.pre_generate_questions_for_document(
        subject=subject_name,
        source_file=source_file,
        count=safe_count,
        force_refresh=True,
        replace_existing=False,
    )

    qb_count_query = db.query(models.QuestionBank).filter(
        models.QuestionBank.source_file == source_file,
    )
    if subject_id:
        qb_count_query = qb_count_query.filter(models.QuestionBank.subject_id == subject_id)
    else:
        qb_count_query = qb_count_query.filter(models.QuestionBank.subject == subject_name)

    final_rows = qb_count_query.order_by(models.QuestionBank.id.desc()).all()
    _cleanup_question_bank_rows(final_rows, db)
    final_count = len(final_rows)
    logger.info(
        "append_question_bank done doc_id=%s source_file=%s previous_count=%s generated_count=%s final_count=%s duration_ms=%s",
        doc_id,
        source_file,
        existing_count,
        len(generated),
        final_count,
        round((time.perf_counter() - started) * 1000, 2),
    )

    return {
        "doc_id": doc_id,
        "subject": subject_name,
        "source_file": source_file,
        "generated_count": int(len(generated)),
        "final_count": int(final_count),
        "previous_count": int(existing_count),
        "requested_count": safe_count,
        "generated_questions": [
            {
                "id": item.get("id"),
                "content": _sanitize_question_text(item.get("content", "")),
                "options": _normalize_question_options(item.get("options", [])),
                "correct_answer": item.get("correct_answer"),
                "explanation": _sanitize_question_text(item.get("explanation", "")),
                "difficulty": "Basic",
                "subject": subject_name,
                "source_file": source_file,
            }
            for item in generated
        ],
    }

# ==========================================
# 3. API CHO HỌC SINH: LẤY DANH SÁCH TÀI LIỆU
# ==========================================
@router.get("/student/{user_id}")
def get_student_documents(user_id: int, subject: Optional[str] = None, db: Session = Depends(get_db)):
    """Lấy danh sách tài liệu dành cho học sinh dựa vào các lớp đã tham gia"""
    user = db.query(models.User).filter(models.User.id == user_id).first()
    if not user:
        raise HTTPException(status_code=404, detail="Không tìm thấy học sinh")

    # Lấy danh sách ID các lớp mà học sinh đang tham gia
    enrolled_class_ids = [c.id for c in getattr(user, 'enrolled_classes', [])]
    
    if not enrolled_class_ids:
        return []

    # Chỉ lấy tài liệu thuộc các lớp đã tham gia VÀ được giáo viên bật hiển thị.
    query = db.query(models.Document, models.DocumentPublication).join(
        models.DocumentPublication,
        models.DocumentPublication.doc_id == models.Document.id,
    ).filter(
        models.Document.class_id.in_(enrolled_class_ids),
        models.DocumentPublication.is_visible_to_students == True,
    )
    
    # Nếu có chọn môn cụ thể thì lọc thêm theo môn
    if subject and subject != "Tất cả":
        query = query.filter(models.Document.subject == subject)
        
    docs = query.order_by(models.Document.id.desc()).all()
    
    result = []
    for doc, publication in docs:
        time_str = (doc.upload_time + timedelta(hours=7)).strftime("%H:%M - %d/%m/%Y") if doc.upload_time else "Chưa xác định"
        result.append({
            "id": doc.id,
            "filename": doc.filename,
            "subject": doc.subject,
            "class_id": doc.class_id,
            "upload_time": time_str,
            "is_visible_to_students": publication.is_visible_to_students,
        })
        
    return result

# ==========================================
# 4. API CHUNG: TẢI FILE VỀ MÁY (DOWNLOAD)
# ==========================================
@router.get("/download/{doc_id}")
def download_document(doc_id: int, db: Session = Depends(get_db)):
    """API cho phép tải file vật lý về thiết bị"""
    doc = db.query(models.Document).filter(models.Document.id == doc_id).first()
    if not doc:
        raise HTTPException(status_code=404, detail="Tài liệu không tồn tại")

    publication = db.query(models.DocumentPublication).filter(models.DocumentPublication.doc_id == doc_id).first()
    if not publication or not publication.is_visible_to_students:
        raise HTTPException(status_code=403, detail="Tài liệu này chưa được giáo viên cho phép hiển thị")
        
    # Trỏ đúng vào thư mục temp_uploads theo logic lưu file
    file_path = _resolve_document_file_path(doc.filename)
        
    return FileResponse(
        path=str(file_path), 
        filename=doc.filename,
        media_type='application/octet-stream' # Ép trình duyệt tự động tải file xuống
    )


@router.get("/view/{doc_id}")
def view_document_inline(doc_id: int, user_id: Optional[int] = None, db: Session = Depends(get_db)):
    """API hiển thị file trực tiếp (inline) để nhúng preview PDF/PPT trong frontend."""
    doc = db.query(models.Document).filter(models.Document.id == doc_id).first()
    if not doc:
        raise HTTPException(status_code=404, detail="Tài liệu không tồn tại")

    if not _can_student_access_document(db, doc, user_id):
        raise HTTPException(status_code=403, detail="Tài liệu này chưa được giáo viên cho phép hiển thị")

    file_path = _resolve_document_file_path(doc.filename)

    ext = os.path.splitext(doc.filename)[1].lower()
    media_map = {
        ".pdf": "application/pdf",
        ".ppt": "application/vnd.ms-powerpoint",
        ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        ".txt": "text/plain; charset=utf-8",
        ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    }
    media_type = media_map.get(ext, "application/octet-stream")

    safe_filename = re.sub(r'[^A-Za-z0-9._-]+', '_', doc.filename or 'document')

    def iter_file_bytes(path: Path, chunk_size: int = 1024 * 1024):
        with path.open("rb") as stream:
            while True:
                chunk = stream.read(chunk_size)
                if not chunk:
                    break
                yield chunk

    return StreamingResponse(
        iter_file_bytes(file_path),
        media_type=media_type,
        headers={
            "Content-Disposition": f"inline; filename=\"{safe_filename}\"",
            "Accept-Ranges": "bytes",
            "X-Doc-View-Version": "v4_stream_unicode_fix",
        },
    )


@router.get("/preview/{doc_id}")
def preview_document(doc_id: int, user_id: Optional[int] = None, db: Session = Depends(get_db)):
    """Trích xuất nội dung text để preview trực tiếp trong ứng dụng."""
    doc = db.query(models.Document).filter(models.Document.id == doc_id).first()
    if not doc:
        raise HTTPException(status_code=404, detail="Tài liệu không tồn tại")

    if not _can_student_access_document(db, doc, user_id):
        raise HTTPException(status_code=403, detail="Tài liệu này chưa được giáo viên cho phép hiển thị")

    cache_key = _preview_cache_key(doc.id, doc.filename)
    cached_payload = _get_cached_preview(cache_key)
    if cached_payload:
        return cached_payload

    file_ext = os.path.splitext(doc.filename)[1].lower()

    try:
        fast_segments = _extract_preview_from_vector(doc.filename)
        if fast_segments and file_ext in {".pdf", ".pptx", ".docx", ".txt"}:
            payload = {
                "doc_id": doc.id,
                "filename": doc.filename,
                "file_type": "vector_fallback",
                "segments": fast_segments,
            }
            _set_cached_preview(cache_key, payload)
            return payload

        file_path = _resolve_document_file_path(doc.filename)
        parsed = _extract_preview_segments(str(file_path), doc.filename)
        segments = parsed.get("segments", [])
        if not segments:
            segments = fast_segments

        payload = {
            "doc_id": doc.id,
            "filename": doc.filename,
            "file_type": parsed.get("type", "unknown"),
            "segments": segments,
        }
        _set_cached_preview(cache_key, payload)
        return payload
    except Exception as e:
        fallback_segments = _extract_preview_from_vector(doc.filename)
        if fallback_segments:
            payload = {
                "doc_id": doc.id,
                "filename": doc.filename,
                "file_type": "vector_fallback",
                "segments": fallback_segments,
            }
            _set_cached_preview(cache_key, payload)
            return payload
        logger.exception("preview_document failed doc_id=%s filename=%s", doc.id, doc.filename)
        raise HTTPException(status_code=500, detail=f"Không thể preview tài liệu: {str(e)}")


# ==========================================
# 5. API LẤY DANH SÁCH TÀI LIỆU THEO MÔN HỌC (DÀNH CHO TEACHER)
# ==========================================
@router.get("/by-subject/{subject_id}")
def get_documents_by_subject(
    subject_id: int,
    teacher_id: Optional[int] = None,
    class_id: Optional[int] = None,
    db: Session = Depends(get_db),
):
    """Lấy danh sách tất cả tài liệu của một môn học (cho giáo viên quản lý)"""
    subject = db.query(models.Subject).filter(models.Subject.id == subject_id).first()
    if not subject:
        raise HTTPException(status_code=404, detail="Không tìm thấy môn học")

    query = db.query(models.Document).filter(models.Document.subject_id == subject_id)
    if class_id is not None:
        query = query.filter(models.Document.class_id == class_id)
    elif teacher_id is not None:
        query = query.filter(models.Document.teacher_id == teacher_id)

    docs = query.order_by(models.Document.id.desc()).all()
    
    result = []
    for doc in docs:
        time_str = (doc.upload_time.strftime("%H:%M - %d/%m/%Y") if doc.upload_time else "Chưa xác định")
        result.append({
            "id": doc.id,
            "title": doc.title,
            "filename": doc.filename,
            "subject": subject.name,
            "subject_id": doc.subject_id,
            "upload_time": time_str,
            "class_id": doc.class_id,
            "teacher_id": doc.teacher_id,
        })
    
    return result
